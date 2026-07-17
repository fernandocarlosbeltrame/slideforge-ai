from pathlib import Path
from typing import List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from .parser import Block
from .theme import Theme

class PresentationGenerator:
    def __init__(self, theme: Optional[Theme] = None):
        self.theme = theme or Theme()
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.slide_no = 0
        self.banner_path: Optional[str] = None
        self.logo_path: Optional[str] = None

    def _add_branding(self, slide, title: str):
        self.slide_no += 1
        # header
        if self.banner_path and Path(self.banner_path).exists():
            slide.shapes.add_picture(self.banner_path, 0, 0, width=self.prs.slide_width, height=Inches(0.62))
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(0.62))
            shape.fill.solid(); shape.fill.fore_color.rgb = self.theme.primary; shape.line.fill.background()
        # title
        tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.78), Inches(12.1), Inches(0.55))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.name = self.theme.font
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = self.theme.dark
        # footer
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.12), self.prs.slide_width, Inches(0.38))
        line.fill.solid(); line.fill.fore_color.rgb = self.theme.primary; line.line.fill.background()
        ft = slide.shapes.add_textbox(Inches(0.45), Inches(7.16), Inches(12.3), Inches(0.2))
        p = ft.text_frame.paragraphs[0]
        p.text = f'Reforma Tributária do Consumo  |  {self.slide_no}'
        p.font.name = self.theme.font; p.font.size = Pt(9); p.font.color.rgb = self.theme.white
        p.alignment = PP_ALIGN.RIGHT

    def add_cover(self, title: str, subtitle: str = ''):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        if self.banner_path and Path(self.banner_path).exists():
            slide.shapes.add_picture(self.banner_path, 0, 0, width=self.prs.slide_width, height=Inches(0.9))
        bg = slide.background.fill
        bg.solid(); bg.fore_color.rgb = self.theme.light
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.45), Inches(11.7), Inches(4.6))
        box.fill.solid(); box.fill.fore_color.rgb = self.theme.white; box.line.color.rgb = self.theme.primary
        tf = box.text_frame; tf.clear(); tf.margin_left = Inches(0.5); tf.margin_right = Inches(0.5)
        p = tf.paragraphs[0]; p.text = title; p.font.name=self.theme.font; p.font.size=Pt(34); p.font.bold=True; p.font.color.rgb=self.theme.primary; p.alignment=PP_ALIGN.CENTER
        p2 = tf.add_paragraph(); p2.text = subtitle; p2.font.name=self.theme.font; p2.font.size=Pt(19); p2.font.color.rgb=self.theme.dark; p2.alignment=PP_ALIGN.CENTER; p2.space_before=Pt(18)
        self.slide_no = 1
        f = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.3))
        p = f.text_frame.paragraphs[0]; p.text = 'Material de apresentação corporativa'; p.alignment=PP_ALIGN.CENTER; p.font.name=self.theme.font; p.font.size=Pt(11); p.font.color.rgb=self.theme.dark

    def add_bullet_slide(self, title: str, items: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_branding(slide, title)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.48), Inches(12.0), Inches(5.35))
        card.fill.solid(); card.fill.fore_color.rgb = self.theme.white; card.line.color.rgb = RGBColor(210,220,230)
        tf = card.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.45); tf.margin_right = Inches(0.35); tf.margin_top = Inches(0.25)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.name = self.theme.font
            p.font.size = Pt(18 if len(items) <= 6 else 15)
            p.font.color.rgb = self.theme.dark
            p.space_after = Pt(7)
            p.text = '• ' + p.text

    def add_before_after(self, title: str, before: List[str], after: List[str]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_branding(slide, title)
        for x, heading, items, fill in [(.65,'ANTES',before,RGBColor(252,239,239)),(6.82,'DEPOIS',after,RGBColor(235,248,239))]:
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.48), Inches(5.85), Inches(5.25))
            box.fill.solid(); box.fill.fore_color.rgb=fill; box.line.color.rgb=self.theme.primary
            tf=box.text_frame; tf.clear(); tf.margin_left=Inches(.3); tf.margin_right=Inches(.3)
            p=tf.paragraphs[0]; p.text=heading; p.font.name=self.theme.font; p.font.size=Pt(23); p.font.bold=True; p.font.color.rgb=self.theme.primary; p.alignment=PP_ALIGN.CENTER
            for item in items:
                p=tf.add_paragraph(); p.text='• '+item; p.font.name=self.theme.font; p.font.size=Pt(14); p.font.color.rgb=self.theme.dark; p.space_after=Pt(4)

    def add_timeline(self, title: str, events: List[tuple]):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_branding(slide, title)
        y=2.65
        x0=.75
        width=11.8
        line=slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(y), Inches(width), Inches(.08))
        line.fill.solid(); line.fill.fore_color.rgb=self.theme.secondary; line.line.fill.background()
        step=width/max(1,len(events)-1)
        for i,(year,text) in enumerate(events):
            x=x0+i*step
            dot=slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y-.16), Inches(.38), Inches(.38))
            dot.fill.solid(); dot.fill.fore_color.rgb=self.theme.primary; dot.line.fill.background()
            tb=slide.shapes.add_textbox(Inches(x-.35), Inches(1.55 if i%2==0 else 3.05), Inches(2.1), Inches(1.3))
            tf=tb.text_frame; tf.word_wrap=True
            p=tf.paragraphs[0]; p.text=year; p.font.name=self.theme.font; p.font.bold=True; p.font.size=Pt(18); p.font.color.rgb=self.theme.primary; p.alignment=PP_ALIGN.CENTER
            p2=tf.add_paragraph(); p2.text=text; p2.font.name=self.theme.font; p2.font.size=Pt(12); p2.font.color.rgb=self.theme.dark; p2.alignment=PP_ALIGN.CENTER

    def generate(self, blocks: List[Block], output: str, banner: Optional[str] = None, logo: Optional[str] = None):
        self.banner_path=banner
        self.logo_path=logo
        self.add_cover('Reforma Tributária do Consumo', 'O que muda?')

        # faithful grouping: headings start new logical sections; text split into readable slides.
        current_title='Conteúdo'
        buffer=[]
        sections=[]
        for b in blocks:
            if b.kind in {'title','heading'} or (len(b.text) < 55 and b.text.endswith(':')):
                if buffer:
                    sections.append((current_title, buffer)); buffer=[]
                current_title=b.text.rstrip(':')
            else:
                buffer.append(b.text)
        if buffer:
            sections.append((current_title, buffer))

        for title, items in sections:
            low=title.lower()
            joined=' '.join(items).lower()
            if 'antes x depois' in low or ('antes' in joined and 'depois' in joined and len(items)>8):
                before=[]; after=[]; side='before'
                for it in items:
                    t=it.strip()
                    if 'depois' in t.lower(): side='after'; continue
                    if 'antes' in t.lower(): side='before'; continue
                    (before if side=='before' else after).append(t)
                if before and after:
                    self.add_before_after(title,before,after); continue
            if title.lower()=='transição':
                # preserve detailed transition in normal slides plus a visual timeline
                self.add_timeline('Linha do tempo da transição',[
                    ('2026','Ano teste da CBS e do IBS'),('2027-2028','CBS, fim de PIS/Cofins e início do IS'),('2029-2032','Transição gradual de ICMS/ISS para IBS'),('2033','Vigência integral do novo modelo')])
            # chunk content without dropping anything
            chunk=[]; chars=0
            part=1
            for it in items:
                if chunk and (len(chunk)>=6 or chars+len(it)>850):
                    suffix=f' — Parte {part}' if len(items)>len(chunk) else ''
                    self.add_bullet_slide(title+suffix,chunk)
                    part+=1; chunk=[]; chars=0
                chunk.append(it); chars+=len(it)
            if chunk:
                suffix=f' — Parte {part}' if part>1 else ''
                self.add_bullet_slide(title+suffix,chunk)

        self.prs.save(output)
        return output
