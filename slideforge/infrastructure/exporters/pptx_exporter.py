from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt
from slideforge.application.services.image_fit_calculator import ImageFitCalculator
from slideforge.application.services.slide_geometry import SlideGeometry
from slideforge.application.services.typography_fitter import TypographyFitter
from slideforge.domain.entities.presentation_plan import PresentationPlan
from slideforge.domain.entities.slide_plan import SlidePlan
from slideforge.domain.enums.slide_layout_type import SlideLayoutType
from slideforge.domain.value_objects.bounding_box import BoundingBox
from slideforge.infrastructure.exporters.pptx.components import BulletList, CardGrid, ComparisonComponent, Footer, ImagePanel, TableComponent, TimelineComponent, TitleBlock
from slideforge.infrastructure.exporters.pptx.grid import GridGap, SlideGrid
from slideforge.infrastructure.exporters.pptx.theme_adapter import ThemeAdapter
from slideforge.theme import Theme, get_theme


class PPTXExporter:
    def __init__(self, theme: Theme | None = None):
        self.theme = theme or get_theme("corporate_blue")
        self.theme_adapter = ThemeAdapter(self.theme)
        self.geometry = SlideGeometry()
        self.image_fit = ImageFitCalculator()
        self.text_fit = TypographyFitter()
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.geometry.slide_width)
        self.prs.slide_height = Inches(self.geometry.slide_height)
        self.banner_path: str | None = None
        self.logo_path: str | None = None
        self.deck_title = "SlideForge AI"
        self.title_block = TitleBlock(self.theme_adapter)
        self.footer = Footer(self.theme_adapter)
        self.bullets = BulletList(self.theme_adapter)
        self.cards = CardGrid(self.theme_adapter)
        self.comparison = ComparisonComponent(self.theme_adapter)
        self.timeline = TimelineComponent(self.theme_adapter)
        self.image_panel = ImagePanel(self.theme_adapter)
        self.table_component = TableComponent(self.theme_adapter)

    def export(self, plan: PresentationPlan, output_path: Path, *, banner_path: str | None = None, logo_path: str | None = None) -> Path:
        self.banner_path = banner_path
        self.logo_path = logo_path or banner_path
        self.deck_title = plan.title
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.geometry.slide_width)
        self.prs.slide_height = Inches(self.geometry.slide_height)
        for slide_plan in plan.slides:
            self._render_slide(slide_plan)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output_path)
        return output_path

    def _render_slide(self, slide_plan: SlidePlan) -> None:
        if slide_plan.layout_type == SlideLayoutType.COVER:
            self._cover(slide_plan)
        elif slide_plan.layout_type in {SlideLayoutType.CARDS, SlideLayoutType.PROCESS_FLOW}:
            self._cards(slide_plan)
        elif slide_plan.layout_type == SlideLayoutType.TIMELINE:
            self._timeline(slide_plan)
        elif slide_plan.layout_type == SlideLayoutType.COMPARISON:
            self._comparison(slide_plan)
        elif slide_plan.layout_type == SlideLayoutType.IMAGE:
            self._image(slide_plan)
        elif slide_plan.layout_type == SlideLayoutType.IMAGE_TEXT:
            self._image_text(slide_plan)
        elif slide_plan.layout_type == SlideLayoutType.TABLE:
            self._table(slide_plan)
        elif slide_plan.layout_type == SlideLayoutType.CALLOUT:
            self._callout(slide_plan)
        else:
            self._bullets(slide_plan)

    def _blank_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self.theme.light
        return slide

    def _branding(self, slide, title: str, number: int) -> None:
        self._header(slide)
        self._title(slide, title)
        self._footer(slide, number)

    def _header(self, slide) -> None:
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(self.geometry.header_height))
        header.fill.solid(); header.fill.fore_color.rgb = self.theme.header_background; header.line.fill.background()
        if self.banner_path and Path(self.banner_path).exists():
            try:
                with Image.open(self.banner_path) as img:
                    fit = self.image_fit.calculate(img.width, img.height, BoundingBox(4.95, 0.04, 3.45, 0.70), "contain")
                slide.shapes.add_picture(self.banner_path, Inches(fit.x), Inches(fit.y), width=Inches(fit.width), height=Inches(fit.height))
            except Exception:
                pass
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(self.geometry.header_height - 0.035), self.prs.slide_width, Inches(0.035))
        border.fill.solid(); border.fill.fore_color.rgb = self.theme.header_border; border.line.fill.background()

    def _footer(self, slide, number: int) -> None:
        self.footer.render(slide, self.geometry.footer_box, self.deck_title, number)

    def _title(self, slide, title: str) -> None:
        self.title_block.render(slide, self.geometry.title_box, title, "slide_title")

    def _content_panel(self, slide, box: BoundingBox | None = None):
        box = box or self.geometry.safe_content_box
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box.x), Inches(box.y), Inches(box.width), Inches(box.height))
        panel.fill.solid(); panel.fill.fore_color.rgb = self.theme.white; panel.line.color.rgb = RGBColor(190, 205, 225)
        return panel

    def _cover(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._header(slide)
        box = BoundingBox(1.05, 1.78, 11.25, 3.85)
        panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box.x), Inches(box.y), Inches(box.width), Inches(box.height))
        panel.fill.solid(); panel.fill.fore_color.rgb = self.theme.white; panel.line.color.rgb = self.theme.primary
        tf = panel.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        fit = self.text_fit.fit(plan.title, box.width - 0.6, 1.4, 32, 20, 3)
        p = tf.paragraphs[0]; p.text = plan.title; p.font.name = self.theme.font; p.font.size = Pt(fit.font_size); p.font.bold = True; p.font.color.rgb = self.theme.primary; p.alignment = PP_ALIGN.CENTER
        if plan.subtitle:
            p2 = tf.add_paragraph(); p2.text = plan.subtitle; p2.font.name = self.theme.font; p2.font.size = Pt(16); p2.font.color.rgb = self.theme.dark; p2.alignment = PP_ALIGN.CENTER
        self._footer(slide, plan.sequence)

    def _bullets(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        panel = self._content_panel(slide)
        self.bullets.render(panel, self._text_items(plan), self._body_font_size(self._text_items(plan)))

    def _cards(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        items = self._text_items(plan)[:6]
        if not items:
            return
        rows = 2 if len(items) > 3 else 1
        cols = 3 if len(items) > 4 else (2 if len(items) > 1 else 1)
        grid = SlideGrid(self.geometry.safe_content_box, GridGap(0.28, 0.3))
        row_boxes = grid.rows(rows)
        boxes = []
        for row in row_boxes:
            for col in SlideGrid(row.box, GridGap(0.28, 0.3)).columns(cols):
                boxes.append(col.box)
        self.cards.render(slide, boxes, items)

    def _comparison(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        comp = next((c for c in plan.components if c.get("type") == "comparison"), {"before": [], "after": []})
        left, right = SlideGrid(self.geometry.safe_content_box, GridGap(0.35, 0.28)).split_ratio(0.5)
        self.comparison.render(slide, left, right, comp.get("before", []), comp.get("after", []))

    def _timeline(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        self.timeline.render(slide, self.geometry.safe_content_box, self._text_items(plan))

    def _image(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        image_component = next((c for c in plan.components if c.get("type") == "image" and c.get("path")), None)
        if image_component and self.image_panel.render(slide, image_component["path"], self.geometry.safe_content_box, image_component.get("fit_mode", "contain")):
            return
        self._bullets(plan)

    def _image_text(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        left, right = SlideGrid(self.geometry.safe_content_box, GridGap(0.35, 0.28)).split_ratio(0.45)
        panel = self._content_panel(slide, left)
        self.bullets.render(panel, self._text_items(plan), 12)
        image_component = next((c for c in plan.components if c.get("type") == "image" and c.get("path")), None)
        if image_component:
            self.image_panel.render(slide, image_component["path"], right, image_component.get("fit_mode", "contain"))

    def _callout(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        box = BoundingBox(self.geometry.margin_left + 0.8, self.geometry.content_top + 0.9, self.geometry.content_width - 1.6, 2.2)
        callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(box.x), Inches(box.y), Inches(box.width), Inches(box.height))
        callout.fill.solid(); callout.fill.fore_color.rgb = self.theme.white; callout.line.color.rgb = self.theme.secondary
        tf = callout.text_frame; tf.clear(); tf.word_wrap = True; tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.margin_left = Inches(0.4); tf.margin_right = Inches(0.4); tf.margin_top = Inches(0.3)
        p = tf.paragraphs[0]
        p.text = "\n".join(self._text_items(plan))
        p.font.name = self.theme.font; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = self.theme.primary; p.alignment = PP_ALIGN.CENTER

    def _table(self, plan: SlidePlan) -> None:
        slide = self._blank_slide(); self._branding(slide, plan.title, plan.sequence)
        rows = next((c.get("rows", []) for c in plan.components if c.get("type") == "table"), [])
        self.table_component.render(slide, self.geometry.safe_content_box, rows)

    def _body_font_size(self, items: list[str]) -> int:
        total = sum(len(item) for item in items)
        if len(items) <= 4 and total < 420:
            return 16
        if len(items) <= 7 and total < 760:
            return 13
        return 10

    @staticmethod
    def _text_items(plan: SlidePlan) -> list[str]:
        items: list[str] = []
        for component in plan.components:
            if component.get("type") == "text":
                items.extend(component.get("items", []))
        return [item for item in items if item]

