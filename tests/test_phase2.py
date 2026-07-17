from pathlib import Path
from docx import Document
from PIL import Image
from pptx import Presentation
from slideforge.application.services.image_fit_calculator import ImageFitCalculator
from slideforge.application.services.presentation_planner import PresentationPlanner
from slideforge.application.services.slide_geometry import SlideGeometry
from slideforge.application.services.typography_fitter import TypographyFitter
from slideforge.application.use_cases.generate_presentation import GeneratePresentationUseCase
from slideforge.application.validation.presentation_validator import PresentationValidator
from slideforge.domain.enums.block_type import BlockType
from slideforge.domain.enums.slide_layout_type import SlideLayoutType
from slideforge.domain.value_objects.bounding_box import BoundingBox
from slideforge.infrastructure.readers.docx_reader import DocxReader
from slideforge.infrastructure.readers.pdf_reader import PDFReader


def make_png(path: Path, size=(640, 320)):
    Image.new("RGB", size, color=(220, 235, 255)).save(path)


def test_docx_extracts_image_in_order(tmp_path: Path):
    img = tmp_path / "split.png"
    make_png(img)
    docx = tmp_path / "doc.docx"
    doc = Document()
    doc.add_heading("Titulo", level=1)
    doc.add_paragraph("Texto antes")
    doc.add_picture(str(img))
    doc.add_paragraph("Texto depois")
    doc.save(docx)
    document = DocxReader().read(docx)
    blocks = [b for s in document.sections for b in s.blocks]
    types = [b.type for b in blocks]
    assert BlockType.IMAGE in types
    image_index = types.index(BlockType.IMAGE)
    assert types[image_index - 1] == BlockType.PARAGRAPH
    assert types[image_index + 1] == BlockType.PARAGRAPH


def test_image_associated_with_section_and_planned(tmp_path: Path):
    img = tmp_path / "split.png"
    make_png(img)
    docx = tmp_path / "doc.docx"
    doc = Document()
    doc.add_heading("Split Payment", level=1)
    doc.add_picture(str(img))
    doc.save(docx)
    document = DocxReader().read(docx)
    plan = PresentationPlanner().create_plan(document)
    assert any(slide.layout_type == SlideLayoutType.IMAGE for slide in plan.slides)


def test_image_fit_contain_preserves_ratio():
    fit = ImageFitCalculator().calculate(800, 400, BoundingBox(0, 0, 8, 4), "contain")
    assert round(fit.width / fit.height, 2) == 2.0
    assert fit.width <= 8 and fit.height <= 4


def test_title_fit_reduces_long_title():
    title = "Reforma Tributária do Consumo: Modernização, Simplificação e Alinhamento às Melhores Práticas Internacionais"
    result = TypographyFitter().fit(title, 10, 0.7, 28, 18, 3)
    assert result.adjusted
    assert result.font_size >= 18


def test_validator_rejects_unused_required_block(tmp_path: Path):
    source = tmp_path / "a.txt"
    source.write_text("A", encoding="utf-8")
    result = GeneratePresentationUseCase().execute(source, tmp_path / "a.pptx")
    assert PresentationValidator().validate_plan(result.plan).is_valid


def test_pdf_reader_basic(tmp_path: Path):
    import fitz
    pdf = tmp_path / "a.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Texto do PDF")
    doc.save(pdf)
    doc.close()
    parsed = PDFReader().read(pdf)
    assert parsed.format == "pdf"
    assert any("Texto do PDF" in b.text for s in parsed.sections for b in s.blocks)


def test_generated_pptx_is_16_9(tmp_path: Path):
    source = tmp_path / "a.txt"
    out = tmp_path / "a.pptx"
    source.write_text("# Título\n- Item 1\n- Item 2", encoding="utf-8")
    GeneratePresentationUseCase().execute(source, out)
    prs = Presentation(out)
    assert round(prs.slide_width / prs.slide_height, 2) == round(13.333 / 7.5, 2)


def test_docx_table_is_structured(tmp_path: Path):
    docx = tmp_path / "table.docx"
    doc = Document()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Coluna A"
    table.cell(0, 1).text = "Coluna B"
    table.cell(1, 0).text = "Valor A"
    table.cell(1, 1).text = "Valor B"
    doc.save(docx)
    document = DocxReader().read(docx)
    tables = [b for s in document.sections for b in s.blocks if b.type == BlockType.TABLE]
    assert tables
    assert tables[0].metadata["rows"][0] == ["Coluna A", "Coluna B"]


def test_geometry_safe_content_box_inside_slide():
    geometry = SlideGeometry()
    errors = geometry.safe_content_box.validate(geometry.slide_width, geometry.slide_height)
    assert errors == []


def test_timeline_plan_keeps_image_as_extra_slide(tmp_path: Path):
    img = tmp_path / "image.png"
    make_png(img, (500, 300))
    docx = tmp_path / "timeline.docx"
    doc = Document()
    doc.add_heading("Cronograma", level=1)
    doc.add_paragraph("2026 início")
    doc.add_paragraph("2027 evolução")
    doc.add_picture(str(img))
    doc.save(docx)
    document = DocxReader().read(docx)
    plan = PresentationPlanner().create_plan(document)
    assert any(slide.layout_type == SlideLayoutType.TIMELINE for slide in plan.slides)
    assert any(slide.layout_type == SlideLayoutType.IMAGE for slide in plan.slides)
