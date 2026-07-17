from pathlib import Path
from pptx import Presentation
from slideforge.application.preview import HTMLPreviewGenerator
from slideforge.application.services.content_auditor import ContentAuditor
from slideforge.application.services.content_density_analyzer import ContentDensityAnalyzer
from slideforge.application.services.presentation_planner import PresentationPlanner
from slideforge.domain.entities.content_block import ContentBlock
from slideforge.domain.entities.document_section import DocumentSection
from slideforge.domain.entities.source_document import SourceDocument
from slideforge.domain.enums.block_type import BlockType
from slideforge.domain.enums.slide_layout_type import SlideLayoutType
from slideforge.domain.services.layout_selector import LayoutSelector
from slideforge.domain.value_objects.bounding_box import BoundingBox
from slideforge.infrastructure.exporters.pptx.grid import GridGap, SlideGrid
from slideforge.infrastructure.exporters.pptx_exporter import PPTXExporter
from slideforge.theme import get_theme


def make_section(title: str, items: list[str]) -> DocumentSection:
    blocks = [ContentBlock(id="t1", type=BlockType.TITLE, text=title, hierarchy_level=1, order=1, source_reference="test")]
    for idx, item in enumerate(items, start=2):
        blocks.append(ContentBlock(id=f"b{idx}", type=BlockType.LIST_ITEM, text=item, hierarchy_level=2, order=idx, source_reference="test"))
    return DocumentSection(id="s1", title=title, order=1, blocks=blocks)


def test_slide_grid_columns_and_ratio_are_stable():
    grid = SlideGrid(BoundingBox(1, 1, 10, 4), GridGap(0.5, 0.25))
    cols = grid.columns(4)
    assert len(cols) == 4
    assert round(cols[0].box.width, 2) == round(cols[1].box.width, 2)
    left, right = grid.split_ratio(0.6)
    assert left.width > right.width
    assert round(left.width + right.width + 0.5, 2) == 10


def test_density_analyzer_classifies_critical_content():
    block = ContentBlock(id="b1", type=BlockType.PARAGRAPH, text="texto " * 400, hierarchy_level=1, order=1, source_reference="test")
    density = ContentDensityAnalyzer().analyze_blocks([block])
    assert density.level == "critical"


def test_layout_decision_keeps_reason_and_composition():
    section = make_section("Antes e Depois", ["Antes: processo manual", "Depois: processo automatizado"])
    decision = LayoutSelector().decide_for_section(section)
    assert decision.layout_type == SlideLayoutType.COMPARISON
    assert decision.reason
    assert decision.visual_composition == "two_column_comparison"


def test_corporate_blue_theme_tokens_available():
    theme = get_theme("corporate_blue")
    assert theme.name == "corporate_blue"
    assert theme.primary is not None
    assert theme.secondary is not None


def test_preview_html_is_generated(tmp_path: Path):
    section = make_section("Pontos", ["Item 1", "Item 2", "Item 3"])
    doc = SourceDocument(id="doc", name="teste.txt", source_path="teste.txt", format="txt", sections=[section])
    plan = PresentationPlanner().create_plan(doc)
    audit = ContentAuditor().audit(plan)
    out = HTMLPreviewGenerator().generate(plan, audit, tmp_path / "preview.html")
    assert out.exists()
    assert "Preview estrutural" in out.read_text(encoding="utf-8")


def test_phase3_pptx_generation_valid(tmp_path: Path):
    section = make_section("Linha do tempo", ["2026 início", "2027 evolução", "2033 vigência integral"])
    doc = SourceDocument(id="doc", name="teste.txt", source_path="teste.txt", format="txt", sections=[section])
    plan = PresentationPlanner().create_plan(doc)
    out = tmp_path / "phase3.pptx"
    PPTXExporter(theme=get_theme("corporate_blue")).export(plan, out)
    prs = Presentation(out)
    assert out.exists()
    assert len(prs.slides) == len(plan.slides)
    assert any(slide.layout_type == SlideLayoutType.TIMELINE for slide in plan.slides)


def test_comparison_overflow_does_not_create_empty_side():
    blocks = [ContentBlock(id="t1", type=BlockType.TITLE, text="Antes e Depois", hierarchy_level=1, order=1, source_reference="test")]
    blocks.append(ContentBlock(id="m1", type=BlockType.PARAGRAPH, text="ANTES", hierarchy_level=2, order=2, source_reference="test"))
    for i in range(8):
        blocks.append(ContentBlock(id=f"a{i}", type=BlockType.LIST_ITEM, text=f"Antes item {i}", hierarchy_level=2, order=3+i, source_reference="test"))
    blocks.append(ContentBlock(id="m2", type=BlockType.PARAGRAPH, text="DEPOIS", hierarchy_level=2, order=20, source_reference="test"))
    blocks.append(ContentBlock(id="d1", type=BlockType.LIST_ITEM, text="Depois item único", hierarchy_level=2, order=21, source_reference="test"))
    section = DocumentSection(id="s", title="Antes e Depois", order=1, blocks=blocks)
    doc = SourceDocument(id="doc", name="teste.txt", source_path="teste.txt", format="txt", sections=[section])
    plan = PresentationPlanner().create_plan(doc)
    comparison_slides = [slide for slide in plan.slides if slide.layout_type == SlideLayoutType.COMPARISON]
    assert comparison_slides
    for slide in comparison_slides:
        comp = next(component for component in slide.components if component.get("type") == "comparison")
        assert comp.get("before")
        assert comp.get("after")
    assert any("comparison_overflow" in slide.render_notes for slide in plan.slides)


def test_comparison_marker_with_emoji_and_remainder_keeps_content():
    blocks = [
        ContentBlock(id="t1", type=BlockType.TITLE, text="Antes x Depois", hierarchy_level=1, order=1, source_reference="test"),
        ContentBlock(id="a0", type=BlockType.PARAGRAPH, text="🔴 ANTES\nSistema atual complexo", hierarchy_level=2, order=2, source_reference="test"),
        ContentBlock(id="a1", type=BlockType.LIST_ITEM, text="Muitas obrigações acessórias", hierarchy_level=2, order=3, source_reference="test"),
        ContentBlock(id="d0", type=BlockType.PARAGRAPH, text="🟢 DEPOIS\nSistema nacional automatizado", hierarchy_level=2, order=4, source_reference="test"),
        ContentBlock(id="d1", type=BlockType.LIST_ITEM, text="Menos redundância", hierarchy_level=2, order=5, source_reference="test"),
    ]
    section = DocumentSection(id="s", title="Antes x Depois", order=1, blocks=blocks)
    doc = SourceDocument(id="doc", name="teste.txt", source_path="teste.txt", format="txt", sections=[section])
    plan = PresentationPlanner().create_plan(doc)
    comp_slide = next(slide for slide in plan.slides if slide.layout_type == SlideLayoutType.COMPARISON)
    comp = next(component for component in comp_slide.components if component.get("type") == "comparison")
    assert "Sistema atual complexo" in comp["before"]
    assert "Sistema nacional automatizado" in comp["after"]
