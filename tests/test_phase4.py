import json
from pathlib import Path
from docx import Document
from pptx import Presentation
from slideforge.application.services.content_auditor import ContentAuditor
from slideforge.application.services.presentation_planner import PresentationPlanner
from slideforge.application.use_cases.publish_presentation import PublishPresentationUseCase
from slideforge.domain.entities.content_block import ContentBlock
from slideforge.domain.entities.document_section import DocumentSection
from slideforge.domain.entities.source_document import SourceDocument
from slideforge.domain.enums.block_type import BlockType
from slideforge.infrastructure.assets import AssetManager
from slideforge.infrastructure.renderers import DocxRenderer, HtmlRenderer, JsonRenderer, MarkdownRenderer, PdfRenderer, PptxRenderer
from slideforge.infrastructure.themes import ThemeRegistry


def sample_plan():
    blocks = [
        ContentBlock(id="t", type=BlockType.TITLE, text="Apresentação Teste", hierarchy_level=1, order=1, source_reference="test"),
        ContentBlock(id="b1", type=BlockType.LIST_ITEM, text="Primeiro ponto", hierarchy_level=2, order=2, source_reference="test"),
        ContentBlock(id="b2", type=BlockType.LIST_ITEM, text="Segundo ponto", hierarchy_level=2, order=3, source_reference="test"),
    ]
    section = DocumentSection(id="s", title="Apresentação Teste", order=1, blocks=blocks)
    doc = SourceDocument(id="doc", name="teste.txt", source_path="teste.txt", format="txt", sections=[section])
    plan = PresentationPlanner().create_plan(doc)
    return plan, ContentAuditor().audit(plan)


def test_asset_manager_registers_and_locates_banner(tmp_path: Path):
    banner = tmp_path / "banner.png"
    banner.write_bytes(b"fake")
    manager = AssetManager()
    manager.register_banner(banner)
    assert manager.locate("banner") == banner


def test_theme_registry_exposes_corporate_blue():
    registry = ThemeRegistry()
    assert "corporate_blue" in registry.names()
    assert registry.get("corporate_blue").base.name == "corporate_blue"


def test_json_renderer_preserves_traceability(tmp_path: Path):
    plan, audit = sample_plan()
    out = JsonRenderer().render(plan, tmp_path / "out.json", audit=audit)
    data = json.loads(out.read_text(encoding="utf-8"))
    assert data["slides"]
    assert data["slides"][0]["source_block_ids"]
    assert "audit" in data


def test_markdown_html_docx_pdf_renderers_create_files(tmp_path: Path):
    plan, audit = sample_plan()
    outputs = [
        MarkdownRenderer().render(plan, tmp_path / "out.md", audit=audit),
        HtmlRenderer().render(plan, tmp_path / "out.html", audit=audit),
        DocxRenderer().render(plan, tmp_path / "out.docx", audit=audit),
        PdfRenderer().render(plan, tmp_path / "out.pdf", audit=audit),
    ]
    for out in outputs:
        assert out.exists()
        assert out.stat().st_size > 0
    assert "Primeiro ponto" in (tmp_path / "out.md").read_text(encoding="utf-8")
    assert "Primeiro ponto" in (tmp_path / "out.html").read_text(encoding="utf-8")
    assert Document(tmp_path / "out.docx").paragraphs


def test_pptx_renderer_uses_same_plan(tmp_path: Path):
    plan, audit = sample_plan()
    out = PptxRenderer().render(plan, tmp_path / "out.pptx", audit=audit)
    prs = Presentation(out)
    assert len(prs.slides) == len(plan.slides)


def test_publish_use_case_exports_selected_formats(tmp_path: Path):
    docx = tmp_path / "entrada.docx"
    doc = Document()
    doc.add_heading("Título", level=1)
    doc.add_paragraph("Primeiro ponto", style=None)
    doc.save(docx)
    result = PublishPresentationUseCase(theme_name="corporate_blue").execute(docx, tmp_path / "saida", formats={"json", "markdown", "html"})
    names = {doc.format_name for doc in result.documents}
    assert names == {"json", "markdown", "html"}
    assert result.audit.unused_block_ids == set()


def test_renderer_content_consistency_between_json_and_markdown(tmp_path: Path):
    plan, audit = sample_plan()
    json_path = JsonRenderer().render(plan, tmp_path / "out.json", audit=audit)
    md_path = MarkdownRenderer().render(plan, tmp_path / "out.md", audit=audit)
    data = json.loads(json_path.read_text(encoding="utf-8"))
    markdown = md_path.read_text(encoding="utf-8")
    assert data["title"] in markdown
    assert "Primeiro ponto" in markdown
