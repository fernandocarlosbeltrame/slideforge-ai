from pathlib import Path
from docx import Document
from pptx import Presentation
from slideforge.application.services.content_auditor import ContentAuditor
from slideforge.application.services.presentation_planner import PresentationPlanner
from slideforge.application.use_cases.generate_presentation import GeneratePresentationUseCase
from slideforge.domain.entities.content_block import ContentBlock
from slideforge.domain.entities.document_section import DocumentSection
from slideforge.domain.services.layout_selector import LayoutSelector
from slideforge.domain.enums.block_type import BlockType
from slideforge.domain.enums.slide_layout_type import SlideLayoutType
from slideforge.infrastructure.readers.docx_reader import DocxReader
from slideforge.infrastructure.readers.txt_reader import TxtReader


def test_txt_reader_creates_blocks(tmp_path: Path):
    path = tmp_path / "entrada.txt"
    path.write_text("Titulo\n- Item um\n- Item dois", encoding="utf-8")
    document = TxtReader().read(path)
    assert document.sections
    assert len(document.sections[0].blocks) == 3
    assert document.sections[0].blocks[1].type == BlockType.LIST_ITEM


def test_docx_reader_creates_blocks(tmp_path: Path):
    path = tmp_path / "entrada.docx"
    doc = Document()
    doc.add_heading("Minha apresentação", level=1)
    doc.add_paragraph("Conteúdo importante.")
    doc.save(path)
    document = DocxReader().read(path)
    assert document.format == "docx"
    assert any(block.text == "Minha apresentação" for section in document.sections for block in section.blocks)
    assert any(block.text == "Conteúdo importante." for section in document.sections for block in section.blocks)


def test_layout_selector_cards():
    section = DocumentSection(id="s1", title="Lista curta", blocks=[
        ContentBlock(id="b1", type=BlockType.LIST_ITEM, text="Item 1"),
        ContentBlock(id="b2", type=BlockType.LIST_ITEM, text="Item 2"),
    ])
    assert LayoutSelector().select_for_section(section) == SlideLayoutType.CARDS


def test_layout_selector_timeline():
    section = DocumentSection(id="s1", title="Transição", blocks=[
        ContentBlock(id="b1", type=BlockType.PARAGRAPH, text="2026 teste, 2027 início, 2033 vigência total"),
    ])
    assert LayoutSelector().select_for_section(section) == SlideLayoutType.TIMELINE


def test_long_content_is_split(tmp_path: Path):
    path = tmp_path / "longo.txt"
    path.write_text(" ".join(["conteúdo"] * 400), encoding="utf-8")
    document = TxtReader().read(path)
    plan = PresentationPlanner().create_plan(document)
    assert len(plan.slides) > 2
    audit = ContentAuditor().audit(plan)
    assert audit.split_block_ids


def test_traceability_marks_used_blocks(tmp_path: Path):
    path = tmp_path / "rastreio.txt"
    path.write_text("- A\n- B\n- C", encoding="utf-8")
    document = TxtReader().read(path)
    plan = PresentationPlanner().create_plan(document)
    audit = ContentAuditor().audit(plan)
    all_blocks = {block.id for section in document.sections for block in section.blocks}
    assert all_blocks <= audit.all_tracked_ids


def test_generate_valid_pptx(tmp_path: Path):
    source = tmp_path / "entrada.txt"
    output = tmp_path / "saida.pptx"
    source.write_text("# Título\n- Item 1\n- Item 2\n- Item 3", encoding="utf-8")
    result = GeneratePresentationUseCase().execute(source, output)
    assert result.output_path.exists()
    prs = Presentation(output)
    assert len(prs.slides) >= 2
